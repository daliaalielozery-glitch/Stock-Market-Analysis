"""
Run this script once to generate the PowerPoint file:
    pip install python-pptx
    python generate_ppt.py
It will create:  Stock_Market_Analysis_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ──────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x0D, 0x2B, 0x55)   # slide backgrounds / titles
MID_BLUE    = RGBColor(0x1F, 0x77, 0xB4)   # accent / underlines
LIGHT_BLUE  = RGBColor(0xD6, 0xE8, 0xF7)   # card fills
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF0, 0xF4, 0xF8)
GREEN       = RGBColor(0x2E, 0xCC, 0x71)
ORANGE      = RGBColor(0xF3, 0x96, 0x14)
PURPLE      = RGBColor(0x8E, 0x44, 0xAD)
TEAL        = RGBColor(0x16, 0xA0, 0x85)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank

# ── Helper functions ─────────────────────────────────────────────────────────

def bg(slide, color):
    """Fill slide background with a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill_color, line_color=None, line_width=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h,
        size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def accent_bar(slide, t=0.92, color=MID_BLUE):
    """Thin horizontal accent bar."""
    rect(slide, 0, t, 13.33, 0.06, color)

def slide_number(slide, n, total, color=WHITE):
    txt(slide, f"{n} / {total}", 12.3, 7.1, 0.9, 0.3,
        size=11, color=color, align=PP_ALIGN.RIGHT)

TOTAL = 14   # total slide count

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_BLUE)
accent_bar(s, 0.0, MID_BLUE)
accent_bar(s, 7.44, MID_BLUE)

# Big circle decoration
c = s.shapes.add_shape(9, Inches(9.8), Inches(-1.2), Inches(5), Inches(5))
c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0x1A, 0x3F, 0x6F)
c.line.fill.background()

c2 = s.shapes.add_shape(9, Inches(10.6), Inches(3.5), Inches(3.5), Inches(3.5))
c2.fill.solid(); c2.fill.fore_color.rgb = RGBColor(0x15, 0x35, 0x60)
c2.line.fill.background()

txt(s, "📈", 0.5, 1.4, 1.2, 1.2, size=42, align=PP_ALIGN.CENTER)
txt(s, "Stock Market Analysis System",
    1.5, 1.3, 9, 1.1, size=38, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
txt(s, "A Web-Based Application Built with Streamlit & Python",
    1.5, 2.5, 9, 0.6, size=19, color=LIGHT_BLUE, align=PP_ALIGN.LEFT, italic=True)

# Divider line
rect(s, 1.5, 3.25, 8, 0.04, MID_BLUE)

info = [("👨‍🏫 Supervisor", "Dr. Khalaf"),
        ("👥 Team Size",    "4 Members"),
        ("🛠️ Level",        "Simple — Beginner")]
for i, (label, val) in enumerate(info):
    x = 1.5 + i * 3.5
    txt(s, label, x, 3.5,  3.2, 0.4, size=12, color=LIGHT_BLUE)
    txt(s, val,   x, 3.85, 3.2, 0.4, size=16, bold=True, color=WHITE)

slide_number(s, 1, TOTAL)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_GRAY)
rect(s, 0, 0, 13.33, 1.1, DARK_BLUE)
txt(s, "📋  Agenda", 0.4, 0.22, 12, 0.7, size=28, bold=True, color=WHITE)
slide_number(s, 2, TOTAL, DARK_BLUE)

items = [
    ("01", "Project Overview & Objectives"),
    ("02", "System Architecture"),
    ("03", "Tools & Technologies"),
    ("04", "Data Flow Diagrams (DFD Level 0 & 1)"),
    ("05", "Use Case & Activity Diagrams"),
    ("06", "Sequence & ER Diagrams"),
    ("07", "Team Task Allocation"),
    ("08", "Live Demo / Screenshots"),
    ("09", "Testing & Results"),
    ("10", "Conclusion"),
]
cols = [items[:5], items[5:]]
for ci, col in enumerate(cols):
    for ri, (num, label) in enumerate(col):
        x = 0.5 + ci * 6.5
        y = 1.4 + ri * 1.1
        r = rect(s, x, y, 6.1, 0.85, WHITE, MID_BLUE, Pt(1))
        r.line.color.rgb = LIGHT_BLUE
        txt(s, num,   x+0.15, y+0.12, 0.7, 0.6, size=20, bold=True, color=MID_BLUE)
        txt(s, label, x+0.85, y+0.18, 5.0, 0.55, size=14, color=DARK_BLUE)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Project Overview
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
rect(s, 0, 0, 13.33, 1.1, DARK_BLUE)
txt(s, "01  Project Overview", 0.4, 0.22, 12, 0.7, size=28, bold=True, color=WHITE)
slide_number(s, 3, TOTAL, DARK_BLUE)

overview_cards = [
    ("🎯", "Goal",        "Build a web app to fetch, display\nand analyze stock market data."),
    ("👤", "Target Users","Students & beginners learning\nstock market concepts."),
    ("🗓️", "Level",       "Simple / Beginner — suitable for\nan intro software course."),
    ("🌐", "Platform",    "Web browser — no installation\nneeded for end users."),
]
for i, (icon, title, body) in enumerate(overview_cards):
    x = 0.35 + i * 3.2
    r = rect(s, x, 1.4, 2.9, 4.5, LIGHT_BLUE)
    txt(s, icon,  x+1.15, 1.55, 0.8, 0.8, size=28, align=PP_ALIGN.CENTER)
    txt(s, title, x+0.15, 2.45, 2.6, 0.5, size=15, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    rect(s, x+0.8, 3.0, 1.3, 0.05, MID_BLUE)
    txt(s, body,  x+0.15, 3.1, 2.6, 1.8, size=12, color=DARK_BLUE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Objectives
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_GRAY)
rect(s, 0, 0, 13.33, 1.1, DARK_BLUE)
txt(s, "01  Objectives & Scope", 0.4, 0.22, 12, 0.7, size=28, bold=True, color=WHITE)
slide_number(s, 4, TOTAL, DARK_BLUE)

objs = [
    ("📡", "Collect real-time & historical stock data via Yahoo Finance API"),
    ("📊", "Visualize stock price trends using interactive charts"),
    ("🖥️", "Build a simple, user-friendly Streamlit web interface"),
    ("⚠️", "Handle invalid inputs and API errors gracefully"),
    ("📋", "Display key stock metrics (price, volume, change %)"),
]
for i, (icon, obj) in enumerate(objs):
    y = 1.4 + i * 1.1
    rect(s, 0.5, y, 12.3, 0.9, WHITE, MID_BLUE, Pt(0.75))
    txt(s, icon, 0.7,  y+0.1, 0.7, 0.7, size=22)
    txt(s, obj,  1.5,  y+0.2, 11.0, 0.55, size=15, color=DARK_BLUE)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Tools & Technologies
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
rect(s, 0, 0, 13.33, 1.1, DARK_BLUE)
txt(s, "03  Tools & Technologies", 0.4, 0.22, 12, 0.7, size=28, bold=True, color=WHITE)
slide_number(s, 5, TOTAL, DARK_BLUE)

tools = [
    ("🐍", "Python 3",     "Core programming\nlanguage",       MID_BLUE),
    ("🌊", "Streamlit",    "Frontend + Backend\nweb framework", GREEN),
    ("📦", "yfinance",     "Yahoo Finance\nstock data API",    ORANGE),
    ("🐼", "Pandas",       "Data cleaning\n& processing",      PURPLE),
    ("📉", "Plotly",       "Interactive charts\n& graphs",      TEAL),
]
for i, (icon, name, desc, color) in enumerate(tools):
    x = 0.5 + i * 2.55
    r = rect(s, x, 1.35, 2.3, 3.6, color)
    txt(s, icon, x+0.7,  1.5,  1.0, 0.9, size=30, align=PP_ALIGN.CENTER)
    txt(s, name, x+0.1,  2.5,  2.1, 0.55, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x+0.6, 3.1, 1.1, 0.05, WHITE)
    txt(s, desc, x+0.1,  3.2,  2.1, 1.0,  size=12, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "All tools are free & open-source — no license needed ✅",
    0.5, 5.2, 12.3, 0.5, size=14, color=DARK_BLUE,
    align=PP_ALIGN.CENTER, italic=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — System Architecture
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_GRAY)
rect(s, 0, 0, 13.33, 1.1, DARK_BLUE)
txt(s, "02  System Architecture", 0.4, 0.22, 12, 0.7, size=28, bold=True, color=WHITE)
slide_number(s, 6, TOTAL, DARK_BLUE)

nodes = [
    (0.4,  3.0, 2.1, 1.1, "👤 User\n(Browser)",        DARK_BLUE),
    (3.2,  3.0, 2.1, 1.1, "🌊 Streamlit\nApp",          MID_BLUE),
    (6.0,  1.7, 2.1, 1.1, "📡 Yahoo Finance\nAPI",       GREEN),
    (6.0,  4.3, 2.1, 1.1, "🐼 Pandas\nProcessing",      ORANGE),
    (9.2,  3.0, 2.1, 1.1, "📊 Charts &\nVisualizations", PURPLE),
    (11.5, 3.0, 1.4, 1.1, "🖥️ Output\nDisplay",         TEAL),
]
for (x, y, w, h, label, color) in nodes:
    r = rect(s, x, y, w, h, color)
    txt(s, label, x+0.05, y+0.1, w-0.1, h-0.15, size=12, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)

arrows = [
    (2.5, 3.55, 3.2, 3.55),
    (5.3, 3.2,  6.0, 2.25),
    (5.3, 3.9,  6.0, 4.85),
    (8.1, 2.25, 9.2, 3.3),
    (8.1, 4.85, 9.2, 3.85),
    (11.3, 3.55, 11.5, 3.55),
]
for (x1, y1, x2, y2) in arrows:
    from pptx.util import Pt as UPt
    conn = s.shapes.add_connector(1,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = MID_BLUE
    conn.line.width = UPt(2)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — DFD Level 0 & 1
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
rect(s, 0, 0, 13.33, 1.1, DARK_BLUE)
txt(s, "04  Data Flow Diagrams", 0.4, 0.22, 12, 0.7, size=28, bold=True, color=WHITE)
slide_number(s, 7, TOTAL, DARK_BLUE)

# DFD 0 panel
rect(s, 0.3, 1.2, 5.9, 5.9, LIGHT_BLUE)
txt(s, "DFD Level 0 — Context", 0.5, 1.3, 5.5, 0.5, size=14, bold=True, color=DARK_BLUE)

dfd0 = [("👤 User", 2.0, 1.9), ("⚙️ Stock Market\nSystem", 2.0, 3.1), ("📡 API", 0.7, 4.3), ("📊 Display", 3.3, 4.3)]
for (label, x, y) in dfd0:
    rect(s, x, y, 2.0, 0.8, DARK_BLUE)
    txt(s, label, x+0.05, y+0.1, 1.9, 0.65, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# arrows
for (x1,y1,x2,y2) in [(3.0,2.7,3.0,3.1),(2.2,3.9,1.5,4.3),(3.8,3.9,4.3,4.3)]:
    conn = s.shapes.add_connector(1, Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    conn.line.color.rgb = MID_BLUE; conn.line.width = Pt(1.5)

# DFD 1 panel
rect(s, 7.1, 1.2, 5.9, 5.9, LIGHT_BLUE)
txt(s, "DFD Level 1 — Detail", 7.3, 1.3, 5.5, 0.5, size=14, bold=True, color=DARK_BLUE)

dfd1_steps = ["👤 User", "Enter Stock Symbol", "Fetch Stock Data", "Process Data (Pandas)", "Generate Charts", "Display Output"]
colors1 = [DARK_BLUE, MID_BLUE, GREEN, ORANGE, PURPLE, TEAL]
for i, (step, color) in enumerate(zip(dfd1_steps, colors1)):
    y = 1.9 + i * 0.82
    rect(s, 8.1, y, 3.9, 0.65, color)
    txt(s, step, 8.15, y+0.08, 3.8, 0.52, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(dfd1_steps)-1:
        conn = s.shapes.add_connector(1, Inches(10.0), Inches(y+0.65), Inches(10.0), Inches(y+0.82))
        conn.line.color.rgb = MID_BLUE; conn.line.width = Pt(1.5)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Use Case & Activity Diagrams
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_GRAY)
rect(s, 0, 0, 13.33, 1.1, DARK_BLUE)
txt(s, "04 & 05  Use Case & Activity Diagrams", 0.4, 0.22, 12, 0.7, size=24, bold=True, color=WHITE)
slide_number(s, 8, TOTAL, DARK_BLUE)

# Use Case
rect(s, 0.3, 1.2, 5.9, 5.9, WHITE, MID_BLUE, Pt(1))
txt(s, "Use Case Diagram", 0.6, 1.3, 5.4, 0.4, size=14, bold=True, color=MID_BLUE)
rect(s, 0.5, 1.8, 1.5, 3.6, DARK_BLUE)
txt(s, "👤\nUser", 0.55, 2.9, 1.4, 0.8, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

use_cases = ["🔍 Search Stock", "📋 View Stock Data", "📈 View Charts", "⬇️ Download CSV"]
for i, uc in enumerate(use_cases):
    y = 2.1 + i * 1.15
    rect(s, 2.3, y, 3.5, 0.75, LIGHT_BLUE, MID_BLUE, Pt(0.75))
    txt(s, uc, 2.4, y+0.12, 3.3, 0.52, size=12, color=DARK_BLUE)
    conn = s.shapes.add_connector(1, Inches(2.0), Inches(y+0.375), Inches(2.3), Inches(y+0.375))
    conn.line.color.rgb = MID_BLUE; conn.line.width = Pt(1)

# Activity
rect(s, 7.1, 1.2, 5.9, 5.9, WHITE, MID_BLUE, Pt(1))
txt(s, "Activity Diagram", 7.4, 1.3, 5.4, 0.4, size=14, bold=True, color=MID_BLUE)

act = [("Start",              DARK_BLUE, 1.85),
       ("Enter Stock Symbol", MID_BLUE,  2.65),
       ("Valid Symbol?",       ORANGE,    3.45),
       ("Fetch & Process",    GREEN,     4.25),
       ("Display Charts",     TEAL,      5.05),
       ("End",                DARK_BLUE, 5.85)]
for label, color, y in act:
    rect(s, 8.4, y, 3.2, 0.6, color)
    txt(s, label, 8.45, y+0.1, 3.1, 0.45, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for i in range(len(act)-1):
    y_from = act[i][2] + 0.6
    y_to   = act[i+1][2]
    conn = s.shapes.add_connector(1, Inches(10.0), Inches(y_from), Inches(10.0), Inches(y_to))
    conn.line.color.rgb = MID_BLUE; conn.line.width = Pt(1.5)

txt(s, "No → Show Error", 11.8, 3.65, 1.4, 0.4, size=10, color=ORANGE)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Sequence Diagram
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
rect(s, 0, 0, 13.33, 1.1, DARK_BLUE)
txt(s, "06  Sequence Diagram", 0.4, 0.22, 12, 0.7, size=28, bold=True, color=WHITE)
slide_number(s, 9, TOTAL, DARK_BLUE)

actors = [("👤 User", 1.3, DARK_BLUE), ("🌊 Streamlit App", 5.2, MID_BLUE), ("📡 Stock API", 9.5, GREEN)]
for (label, x, color) in actors:
    rect(s, x, 1.2, 2.3, 0.7, color)
    txt(s, label, x+0.05, 1.3, 2.2, 0.55, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # lifeline
    conn = s.shapes.add_connector(1, Inches(x+1.15), Inches(1.9), Inches(x+1.15), Inches(6.8))
    conn.line.color.rgb = LIGHT_BLUE; conn.line.width = Pt(1)

msgs = [
    (2.45, 5.2,  2.3, "Enter stock symbol",  MID_BLUE, False),
    (6.35, 9.5,  3.1, "Request data",         GREEN,   False),
    (9.5,  6.35, 3.9, "Return stock data",    GREEN,   True),
    (5.2,  2.45, 4.7, "Display charts & info",MID_BLUE,True),
]
for (x1, x2, y, label, color, dashed) in msgs:
    conn = s.shapes.add_connector(1, Inches(x1+1.15), Inches(y), Inches(x2+1.15), Inches(y))
    conn.line.color.rgb = color
    conn.line.width = Pt(1.5)
    if dashed:
        conn.line.dash_style = 4  # dash
    mid_x = min(x1, x2) + abs(x2-x1)/2
    txt(s, label, mid_x+0.5, y-0.35, abs(x2-x1), 0.35, size=11, color=DARK_BLUE)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — ER Diagram
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_GRAY)
rect(s, 0, 0, 13.33, 1.1, DARK_BLUE)
txt(s, "07  ER Diagram (Database — Simple)", 0.4, 0.22, 12, 0.7, size=26, bold=True, color=WHITE)
slide_number(s, 10, TOTAL, DARK_BLUE)

def er_entity(slide, title, fields, x, y, w=3.5, color=MID_BLUE):
    rect(slide, x, y, w, 0.55, color)
    txt(slide, title, x+0.1, y+0.07, w-0.2, 0.42, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for i, (ftype, fname) in enumerate(fields):
        bg_c = LIGHT_BLUE if i % 2 == 0 else WHITE
        rect(slide, x, y+0.55+i*0.5, w, 0.5, bg_c, LIGHT_BLUE, Pt(0.5))
        txt(slide, ftype, x+0.1, y+0.62+i*0.5, 0.9, 0.38, size=11, color=MID_BLUE, italic=True)
        txt(slide, fname, x+1.1, y+0.62+i*0.5, w-1.2, 0.38, size=11, color=DARK_BLUE)

er_entity(s, "USER",
    [("int","user_id  🔑"),("string","name"),("string","email")],
    0.5, 1.5, color=DARK_BLUE)

er_entity(s, "STOCK",
    [("string","stock_symbol  🔑"),("string","company_name")],
    9.3, 1.5, color=GREEN)

er_entity(s, "SEARCH_HISTORY",
    [("int","search_id  🔑"),("string","stock_symbol"),("date","search_date")],
    4.4, 3.8, color=ORANGE)

# relationship lines
for (x1,y1,x2,y2) in [(4.0,3.5,4.4,4.6),(9.3,2.8,7.9,4.55)]:
    conn = s.shapes.add_connector(1, Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    conn.line.color.rgb = DARK_BLUE; conn.line.width = Pt(1.5)

txt(s, "1 : N", 3.4, 3.9, 1.0, 0.35, size=11, color=MID_BLUE, bold=True)
txt(s, "1 : N", 8.0, 3.7, 1.0, 0.35, size=11, color=GREEN,    bold=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Team Task Allocation
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
rect(s, 0, 0, 13.33, 1.1, DARK_BLUE)
txt(s, "07  Team Task Allocation", 0.4, 0.22, 12, 0.7, size=28, bold=True, color=WHITE)
slide_number(s, 11, TOTAL, DARK_BLUE)

members = [
    ("👑", "Member 1", "Project Lead /\nBackend Dev",
     ["Setup Python project structure","Integrate yfinance API","Handle stock symbol validation","Fetch current + historical data","Error handling (invalid symbols)"],
     MID_BLUE),
    ("🔬", "Member 2", "Data Processing\n& Analysis",
     ["Clean data with Pandas","Extract OHLCV fields","Prepare data for charts","Compute MA indicators"],
     GREEN),
    ("📊", "Member 3", "Visualization\nDeveloper",
     ["Line, Candlestick, Area charts","Volume & daily change charts","Interactive Plotly visuals","Dynamic chart updates"],
     ORANGE),
    ("🖥️", "Member 4", "UI Dev /\nTester",
     ["Build Streamlit interface","Layout sidebar + metrics","Test valid/invalid inputs","Prepare documentation & PPT"],
     PURPLE),
]
for i, (icon, name, role, tasks, color) in enumerate(members):
    x = 0.25 + i * 3.27
    rect(s, x, 1.2, 3.0, 5.8, color)
    txt(s, icon, x+1.0, 1.3, 1.0, 0.7, size=24, align=PP_ALIGN.CENTER)
    txt(s, name, x+0.1, 2.05, 2.8, 0.45, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x+0.5, 2.55, 2.0, 0.05, WHITE)
    txt(s, role, x+0.1, 2.65, 2.8, 0.65, size=11, color=WHITE, align=PP_ALIGN.CENTER, italic=True)
    for j, task in enumerate(tasks):
        txt(s, f"• {task}", x+0.15, 3.45+j*0.68, 2.7, 0.62, size=10, color=WHITE)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Key Features / Screenshots
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_GRAY)
rect(s, 0, 0, 13.33, 1.1, DARK_BLUE)
txt(s, "08  Key Features", 0.4, 0.22, 12, 0.7, size=28, bold=True, color=WHITE)
slide_number(s, 12, TOTAL, DARK_BLUE)

features = [
    ("📡", "Real-Time Data",    "Fetches live prices from\nYahoo Finance API",       MID_BLUE),
    ("📊", "3 Chart Types",     "Line, Candlestick & Area\ncharts with Plotly",       GREEN),
    ("📐", "Key Metrics",       "Current price, open, high,\nlow, volume & % change", ORANGE),
    ("📋", "Data Table",        "Sortable historical data\nwith CSV download",         PURPLE),
    ("🔍", "Stock Search",      "Search any listed stock\nwith instant results",       TEAL),
    ("📈", "Moving Averages",   "5-day & 20-day MA\noverlaid on price chart",         MID_BLUE),
]
for i, (icon, title, desc, color) in enumerate(features):
    col, row = i % 3, i // 3
    x = 0.4 + col * 4.3
    y = 1.4 + row * 2.8
    r = rect(s, x, y, 4.0, 2.4, color)
    txt(s, icon,  x+1.6, y+0.2, 0.9, 0.8, size=26, align=PP_ALIGN.CENTER)
    txt(s, title, x+0.1, y+1.0, 3.8, 0.5, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, desc,  x+0.1, y+1.55, 3.8, 0.8, size=12, color=WHITE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Testing
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
rect(s, 0, 0, 13.33, 1.1, DARK_BLUE)
txt(s, "09  Testing & Results", 0.4, 0.22, 12, 0.7, size=28, bold=True, color=WHITE)
slide_number(s, 13, TOTAL, DARK_BLUE)

tests = [
    ("✅", "Valid Symbol — AAPL",      "App fetched data and displayed charts correctly.",  GREEN),
    ("✅", "Valid Symbol — TSLA",       "All metrics and graphs loaded with correct values.", GREEN),
    ("❌", "Invalid Symbol — XYZABC",  "App showed a clear red error message.",              RGBColor(0xC0, 0x39, 0x2B)),
    ("❌", "Empty Input",              "App prompted the user to enter a symbol.",           RGBColor(0xC0, 0x39, 0x2B)),
    ("✅", "CSV Download",             "Data exported successfully as a .csv file.",         GREEN),
    ("✅", "Period Change",            "Charts updated dynamically when period changed.",    GREEN),
    ("✅", "Chart Type Switch",        "Switching between chart types worked instantly.",    GREEN),
]
for i, (icon, test, result, color) in enumerate(tests):
    y = 1.4 + i * 0.83
    rect(s, 0.4, y, 12.5, 0.72, LIGHT_GRAY, color, Pt(1))
    txt(s, icon,   0.6,  y+0.1, 0.6,  0.52, size=16)
    txt(s, test,   1.2,  y+0.15, 4.5, 0.45, size=12, bold=True, color=DARK_BLUE)
    txt(s, result, 5.8,  y+0.15, 7.0, 0.45, size=12, color=DARK_BLUE)
    rect(s, 12.3, y+0.15, 0.5, 0.45, color)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Conclusion
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_BLUE)

c = s.shapes.add_shape(9, Inches(-1), Inches(-1), Inches(6), Inches(6))
c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0x1A, 0x3F, 0x6F)
c.line.fill.background()

c2 = s.shapes.add_shape(9, Inches(9.5), Inches(3.5), Inches(5), Inches(5))
c2.fill.solid(); c2.fill.fore_color.rgb = RGBColor(0x15, 0x35, 0x60)
c2.line.fill.background()

txt(s, "🎯", 5.8, 0.7, 1.8, 1.2, size=40, align=PP_ALIGN.CENTER)
txt(s, "Conclusion", 2.0, 1.8, 9.3, 1.0, size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
rect(s, 4.0, 2.9, 5.3, 0.06, MID_BLUE)

points = [
    "✅  Successfully built a real-time Stock Market Analysis web app",
    "✅  Integrated Yahoo Finance API for live data fetching",
    "✅  Created interactive Plotly charts (Line, Candlestick, Area)",
    "✅  Clean Streamlit UI with error handling",
    "✅  Full team collaboration with defined roles",
]
for i, p in enumerate(points):
    txt(s, p, 1.8, 3.2+i*0.72, 9.7, 0.65, size=14, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "Thank You  🙏", 3.0, 6.8, 7.3, 0.55,
    size=20, bold=True, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

slide_number(s, 14, TOTAL)

# ── Save ────────────────────────────────────────────────────────────────────
fname = "Stock_Market_Analysis_Presentation.pptx"
prs.save(fname)
print(f"✅  Presentation saved as:  {fname}")
print(f"    Slides: {TOTAL}  |  Open with Microsoft PowerPoint or LibreOffice Impress")
